"""Generate SudokuOCR PowerPoint presentation."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGURES = Path(__file__).parent / "figures"
OUT = FIGURES / "SudokuOCR_presentation.pptx"

# ── Colors ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
RED = RGBColor(0xE8, 0x45, 0x3C)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_bg(slide, color=WHITE):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with single-run formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=18,
                color=BLACK, spacing=Pt(6), font_name="Calibri"):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        # Remove existing buNone if present
        for child in list(pPr):
            if child.tag.endswith('}buNone') or child.tag.endswith('}buChar'):
                pPr.remove(child)
        pPr.append(buChar)
    return txBox


def add_image_scaled(slide, img_path, left, top, max_w, max_h):
    """Add image scaled to fit within max_w x max_h, preserving aspect ratio."""
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    target_aspect = max_w / max_h
    if aspect > target_aspect:
        w = max_w
        h = int(max_w / aspect)
    else:
        h = max_h
        w = int(max_h * aspect)
    # Center within the bounding box
    actual_left = left + (max_w - w) // 2
    actual_top = top + (max_h - h) // 2
    slide.shapes.add_picture(str(img_path), actual_left, actual_top, w, h)


def add_accent_bar(slide, left, top, width, height, color=RED):
    """Add a colored accent rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def set_notes(slide, text):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, DARK)

    # Accent bar
    add_accent_bar(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(0.06), RED)

    # Title
    add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.8),
                "SudokuOCR", font_size=54, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.7),
                "Constraint-Driven AI for Perfect Historical Census Digitization",
                font_size=24, color=RGBColor(0xCC, 0xCC, 0xCC),
                alignment=PP_ALIGN.CENTER)

    # Key stat
    add_textbox(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
                "100% accuracy on 16,542 cells across 52 tables",
                font_size=28, bold=True, color=RED,
                alignment=PP_ALIGN.CENTER)

    # Stats row
    stats = [
        ("52", "Tables\nPerfect"),
        ("16,542", "Cells\nVerified"),
        ("26,173", "Constraint\nChecks Passed"),
        ("$0.08", "Per Table\nAverage Cost"),
    ]
    x_start = Inches(1.5)
    for i, (num, label) in enumerate(stats):
        x = x_start + Inches(i * 2.8)
        add_textbox(slide, x, Inches(4.8), Inches(2.4), Inches(0.6),
                    num, font_size=36, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.4), Inches(2.4), Inches(0.7),
                    label, font_size=14, color=RGBColor(0xAA, 0xAA, 0xAA),
                    alignment=PP_ALIGN.CENTER)

    # Author line
    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
                "Chris Ruano  |  HEB 91r  |  Harvard University",
                font_size=16, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "Welcome. I'm going to show you how we achieved something that sounds impossible: "
        "100% cell-level accuracy on historical census tables using a constraint-driven AI approach. "
        "52 tables, over 16,000 cells, zero errors. And the system costs about 8 cents per table.")


def slide_02_problem(prs):
    """The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(6), Inches(0.7),
                "The Problem", font_size=36, bold=True, color=DARK)

    # Left column: text
    bullets = [
        "Historical Indian Census (1872\u20131941): millions of pages of handwritten/printed tables",
        "20+ provinces, 70+ years, thousands of tables of demographic data",
        "Manual transcription: ~$10/table, slow, and still has errors",
        "We found 28 confirmed errors in professionally transcribed ground truth",
    ]
    add_bullets(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(4.0),
                bullets, font_size=18, color=BLACK)

    # Right: example scan
    scan = FIGURES / "example_scan.png"
    if scan.exists():
        add_image_scaled(slide, scan,
                         Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.0))

    # Caption under image
    add_textbox(slide, Inches(6.5), Inches(6.3), Inches(6.3), Inches(0.5),
                "Travancore Eastern Division, 1901 Census",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "The Indian Census from 1872-1941 is one of the largest historical demographic datasets ever collected. "
        "It covers 20+ provinces over 70 years. Each page contains age-by-sex population tables. "
        "Manual transcription costs about $10 per table and is slow. Worse, it's not even reliable -- "
        "we found 28 confirmed errors in our professionally transcribed ground truth files. "
        "Some errors were off by thousands.")


def slide_03_why_hard(prs):
    """Why This Is Hard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(8), Inches(0.7),
                "Why This Is Hard", font_size=36, bold=True, color=DARK)

    bullets = [
        "Tables have 240+ cells with complex structure (sections, age groups, P/M/F columns)",
        "OCR errors are correlated: same digit confusion across an entire page",
        "Even the best AI models plateau at ~92\u201396% \u2014 the \"last mile\" problem",
        "4% error rate \u00d7 240 cells = ~10 wrong cells per table, with no way to know which",
        "Common confusions: 3\u21948, 5\u21946, 0\u21949 (ambiguous in degraded scans)",
    ]
    add_bullets(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4.5),
                bullets, font_size=18, color=BLACK)

    # Right side: stylized table structure
    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.4), Inches(5.3), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("     POPULATION        UNMARRIED        MARRIED\n"
              "  P     M     F     P     M     F     P     M     F\n"
              " \u2500\u2500\u2500\u2500 \u2500\u2500\u2500\u2500 \u2500\u2500\u2500\u2500  \u2500\u2500\u2500\u2500 \u2500\u2500\u2500\u2500 \u2500\u2500\u2500\u2500  \u2500\u2500\u2500\u2500 \u2500\u2500\u2500\u2500 \u2500\u2500\u2500\u2500\n"
              " 175k  85k  90k   174k  85k  89k    1.4k  509   958\n"
              " 236k 115k 121k   235k 114k 120k    1.4k  763   727\n"
              "  ...  ...  ...    ...  ...  ...     ...  ...   ...\n\n"
              "20 age rows \u00d7 4 groups \u00d7 3 columns\n"
              "= 240 cells per table")
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK

    # Highlight box
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.8),
                "The fundamental problem: the AI doesn\u2019t know when it\u2019s wrong.",
                font_size=20, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "A single table can have 240 cells or more. With 4% error rate, that's 10 wrong cells per table. "
        "And you have no way to know WHICH cells are wrong. The errors are systematic -- "
        "the same ambiguous digit gets misread the same way every time. "
        "This is the 'last mile' problem: getting from 96% to 100%.")


def slide_04_baseline(prs):
    """The Baseline -- performance comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                'The Baseline \u2014 "Just Ask the AI"', font_size=36, bold=True, color=DARK)

    # Full-width figure
    fig = FIGURES / "fig1_performance.png"
    if fig.exists():
        add_image_scaled(slide, fig,
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))

    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                "Raw AI extraction is nowhere near sufficient. SudokuOCR: 100%.",
                font_size=18, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "Here's every method we tested on all 52 tables. Traditional OCR tools like Tesseract get around 40-50%. "
        "State-of-the-art LLMs like GPT-4.1 and Claude get around 15-16%. "
        "Mistral's dedicated OCR model reaches 71%. Google Cloud Vision gets 82%. "
        "SudokuOCR achieves 100% -- the only method that gets every single cell correct. "
        "The right panel shows per-table distributions: other methods have huge variance, "
        "while SudokuOCR is a flat line at 100%.")


def slide_05_sudoku_insight(prs):
    """The Key Insight -- Sudoku analogy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(12), Inches(0.7),
                "The Key Insight \u2014 Census Tables Are Sudoku Puzzles",
                font_size=36, bold=True, color=DARK)

    # Sudoku analogy figure
    fig = FIGURES / "fig5_sudoku_analogy.png"
    if fig.exists():
        add_image_scaled(slide, fig,
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.0))

    bullets = [
        "Every cell is constrained by its row and column sums",
        "One wrong digit violates multiple constraints simultaneously",
        "Errors become detectable and correctable \u2014 like Sudoku",
    ]
    add_bullets(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2),
                bullets, font_size=16, color=DARK)

    set_notes(slide,
        "This is the key insight: census tables aren't just grids of independent numbers. "
        "They're constraint systems, like Sudoku puzzles. Every row must satisfy Persons = Males + Females. "
        "Column sums must equal totals. Cross-group sums must be consistent. "
        "A single wrong digit violates multiple constraints. This makes errors DETECTABLE and CORRECTABLE.")


def slide_06_constraints(prs):
    """Five Levels of Constraints."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                "Five Levels of Constraints", font_size=36, bold=True, color=DARK)

    levels = [
        ("L1", "Row Constraint", "Persons = Males + Females (every row)", RED),
        ("L2", "Column Sums", "Age groups must sum to their subtotals and totals", RGBColor(0xE6, 0x5C, 0x00)),
        ("L3", "Cross-Group", "Population = Unmarried + Married + Widowed + Divorced", RGBColor(0xF9, 0xA8, 0x25)),
        ("L4", "Cross-Section", "All Communities = Brahmanic + Other Hindus + Muslims + ...", GREEN),
        ("L5", "Non-Negativity", "All values \u2265 0", BLUE),
    ]

    y = Inches(1.5)
    for tag, name, desc, color in levels:
        # Tag box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.8), y, Inches(0.9), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)

        # Name
        add_textbox(slide, Inches(2.0), y, Inches(2.5), Inches(0.6),
                    name, font_size=20, bold=True, color=DARK)
        # Description
        add_textbox(slide, Inches(4.5), y, Inches(8.3), Inches(0.6),
                    desc, font_size=17, color=BLACK)

        y += Inches(0.85)

    # Summary stat
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0),
                "26,173 total constraint checks across 52 tables (~1.6 per cell)\n"
                "Each cell participates in 3\u20135 independent constraints",
                font_size=20, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "We define 5 levels of constraints. L1: every row, Persons = Males + Females. "
        "L2: vertical sums -- age groups must add up to their subtotals and grand totals. "
        "L3: cross-group -- Population must equal the sum of Unmarried + Married + Widowed + Divorced. "
        "L4: cross-section -- community sections must sum to the All Communities total. "
        "L5: all values must be non-negative. "
        "Across our 52 tables, this gives us over 26,000 constraint checks. "
        "Each cell participates in 3-5 independent constraints, so errors can't hide.")


def slide_07_architecture(prs):
    """The Architecture -- 3 Phases."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                "The Architecture \u2014 3 Phases", font_size=36, bold=True, color=DARK)

    fig = FIGURES / "fig4_architecture.png"
    if fig.exists():
        add_image_scaled(slide, fig,
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8))

    set_notes(slide,
        "The pipeline has 3 phases. Phase 1: Schema Discovery -- one API call to understand the table structure. "
        "What are the column groups? What are the age rows? What's the subtotal hierarchy? "
        "Phase 2: Tailored Extraction -- schema-guided data extraction using a prompt generated from the schema. "
        "Phase 3: Verify and Repair -- check all constraints, and if any fail, run the multi-phase repair cascade. "
        "Total: 2-3 API calls per table for most cases.")


def slide_08_repair(prs):
    """The Repair Cascade."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                "The Repair Cascade", font_size=36, bold=True, color=DARK)

    fig = FIGURES / "fig6_repair_cascade.png"
    if fig.exists():
        add_image_scaled(slide, fig,
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.0))

    # Key takeaway
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
                "59% of fixes require ZERO additional API calls \u2014 pure deductive reasoning",
                font_size=20, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "When constraints fail, we don't just re-ask the AI. We have a multi-phase repair cascade. "
        "Phase A detects Male/Female column swaps -- zero API calls. "
        "Phase B uses 10 deductive strategies to fix digit confusions -- also zero API calls. "
        "These cover common OCR confusions like 3-vs-8, 5-vs-6. "
        "Only if deduction fails do we make additional API calls for targeted rechecks. "
        "59% of all fixes are completely free -- pure math, no additional AI needed. "
        "Like Sudoku: use what you know to deduce what you don't.")


def slide_09_results(prs):
    """Results -- Perfect Accuracy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                "Results \u2014 Perfect Accuracy", font_size=36, bold=True, color=DARK)

    fig = FIGURES / "fig7_scale.png"
    if fig.exists():
        add_image_scaled(slide, fig,
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.0))

    # Headline stats
    stats_text = ("52/52 tables perfect  \u00b7  16,542 cells, 0 errors  \u00b7  "
                  "26,173 constraints, 0 failures  \u00b7  6 regions, 4 census years")
    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                stats_text, font_size=18, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "The results: 52 out of 52 tables achieve 100% cell-level accuracy. "
        "16,542 cells with zero errors. 26,173 constraint checks with zero failures. "
        "This spans 6 regions and 4 census years from 1891 to 1941. "
        "The tables range from simple 2-group tables with 114 cells to complex "
        "3-section, 5-group tables with 500+ cells each.")


def slide_10_cost(prs):
    """How We Compare -- cost vs accuracy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                "How We Compare", font_size=36, bold=True, color=DARK)

    fig = FIGURES / "fig2_cost.png"
    if fig.exists():
        add_image_scaled(slide, fig,
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))

    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
                "SudokuOCR: the only method that achieves verifiable perfection",
                font_size=18, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "This chart compares cost versus accuracy for every method we tested. "
        "Manual transcription costs about $10 per table and still has errors -- we found 28 in our GT files. "
        "Traditional OCR is free but only 40-50% accurate. Raw LLMs cost 3-6 cents but plateau at 15-82%. "
        "SudokuOCR costs about 8 cents per table and achieves 100% with verification. "
        "It's the only dot in the upper-left corner: cheap AND perfect AND self-verifying.")


def slide_11_gt_errors(prs):
    """Finding Errors in Ground Truth."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(12), Inches(0.7),
                'Finding Errors in "Ground Truth"', font_size=36, bold=True, color=DARK)

    # Error summary table as text boxes
    headers = ["GT Source", "Errors Found", "Details"]
    data = [
        ["Travancore.xlsx", "4 cells", "Off by ~1,000 at age 30\u201335"],
        ["Hyderabad.xlsx", "21 cells", "Some off by 6,000+"],
        ["Coorg.xlsx", "3 cells", "Off by 11 at age 50\u201355"],
    ]

    # Table header
    y = Inches(1.5)
    cols = [Inches(1.0), Inches(5.0), Inches(8.0)]
    widths = [Inches(3.5), Inches(2.5), Inches(4.5)]
    for i, h in enumerate(headers):
        add_textbox(slide, cols[i], y, widths[i], Inches(0.5),
                    h, font_size=18, bold=True, color=WHITE)

    # Header background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), y, Inches(11.7), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    # Send to back by moving to position 0
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    # Re-add headers on top
    for i, h in enumerate(headers):
        add_textbox(slide, cols[i], y, widths[i], Inches(0.5),
                    h, font_size=18, bold=True, color=WHITE)

    y += Inches(0.6)
    for row in data:
        for i, cell in enumerate(row):
            clr = RED if i == 1 else BLACK
            add_textbox(slide, cols[i], y, widths[i], Inches(0.45),
                        cell, font_size=17, color=clr, bold=(i == 1))
        y += Inches(0.55)

    # Total
    add_textbox(slide, Inches(1.0), y + Inches(0.2), Inches(11.3), Inches(0.5),
                "28 total human transcription errors found across 3 ground truth files",
                font_size=20, bold=True, color=RED)

    # Key points
    bullets = [
        "In every case, the pipeline's constraint-verified output was correct",
        "Constraints serve as a trust signal \u2014 no external ground truth needed",
        "The system is more reliable than human transcription",
    ]
    add_bullets(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5),
                bullets, font_size=18, color=DARK)

    set_notes(slide,
        "Here's an unexpected result: our pipeline is MORE accurate than the human-transcribed ground truth. "
        "We found 28 confirmed errors in 3 professionally transcribed Excel files. "
        "Hyderabad alone had 21 errors, some off by over 6,000. "
        "In every single case, the pipeline's constraint-verified output matched the original scan correctly. "
        "This demonstrates that constraints provide a trust signal independent of any external ground truth. "
        "You don't need a human to verify the output -- the math does it for you.")


def slide_12_broader_principle(prs):
    """The Broader Principle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)

    # Central message
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
                "AI + Verifiable Constraints = Reliable AI",
                font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(4), Inches(2.8), Inches(5.3), Inches(0.05), RED)

    # Pattern explanation
    add_textbox(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6),
                "The pattern works whenever the domain has internal consistency rules:",
                font_size=20, color=RGBColor(0xCC, 0xCC, 0xCC),
                alignment=PP_ALIGN.CENTER)

    examples = [
        ("Census Tables", "AI reads digits \u2192 arithmetic verifies"),
        ("Financial Statements", "AI extracts numbers \u2192 accounting identities verify"),
        ("Scientific Data", "AI reads measurements \u2192 physical laws verify"),
        ("Code Generation", "AI writes code \u2192 tests verify"),
    ]

    y = Inches(4.0)
    for domain, desc in examples:
        add_textbox(slide, Inches(2.0), y, Inches(3.5), Inches(0.5),
                    domain, font_size=20, bold=True, color=RED,
                    alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(5.8), y, Inches(6.0), Inches(0.5),
                    desc, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB))
        y += Inches(0.6)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
                "Use AI for the hard part (perception) \u2014 use math for the easy part (verification)",
                font_size=18, color=RGBColor(0xAA, 0xAA, 0xAA),
                alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "This is the broader principle. AI plus verifiable constraints equals reliable AI. "
        "The pattern works whenever you have a domain with internal consistency rules. "
        "Census tables have arithmetic. Financial statements have accounting identities. "
        "Scientific data has physical laws. Code has tests. "
        "Use AI for the hard part -- perception, generation -- and use math for the easy part -- verification. "
        "The constraints don't need to directly check correctness. They just need to be CORRELATED with it.")


def slide_13_future(prs):
    """Future Directions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.5), RED)

    add_textbox(slide, Inches(0.75), Inches(0.4), Inches(10), Inches(0.7),
                "Future Directions", font_size=36, bold=True, color=DARK)

    sections = [
        ("Scale", [
            "Full Indian Census corpus: thousands of tables across all provinces",
            "Batch processing pipeline with automatic quality assessment",
            "Estimated cost: ~$77 for 1,000 tables",
        ]),
        ("Extend", [
            "Other structured historical documents: vital statistics, trade records",
            "Multi-page tables with cross-page constraints",
            "Non-English census records (French, Dutch colonial)",
        ]),
        ("Generalize", [
            "Open-source the constraint framework",
            "Apply to any domain with verifiable structure",
            "Financial auditing, scientific data validation, inventory records",
        ]),
    ]

    y = Inches(1.5)
    for heading, items in sections:
        add_textbox(slide, Inches(0.8), y, Inches(3), Inches(0.5),
                    heading, font_size=24, bold=True, color=RED)
        add_bullets(slide, Inches(1.0), y + Inches(0.5), Inches(11.5), Inches(1.5),
                    items, font_size=16, color=BLACK)
        y += Inches(1.9)

    set_notes(slide,
        "Looking ahead, we want to scale this to the full Indian Census corpus -- thousands of tables. "
        "At about 8 cents per table, 1000 tables would cost around $77. "
        "We also want to extend to other structured historical documents: vital statistics, trade records, "
        "and non-English colonial census records. "
        "Most importantly, we want to open-source the constraint framework so it can be applied "
        "to any domain with verifiable structure -- financial auditing, scientific data, inventory management.")


def slide_14_thankyou(prs):
    """Thank You / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0),
                "Thank You", font_size=48, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(4), Inches(2.6), Inches(5.3), Inches(0.05), RED)

    # Summary stats
    summary = [
        "100% cell-level accuracy on 52 historical census tables",
        "16,542 cells verified through 26,173 constraint checks",
        "$0.08 per table \u2014 cheaper and more accurate than manual transcription",
        "Self-verifying: constraint pass rate = quality metric, no GT needed",
    ]
    add_bullets(slide, Inches(2.5), Inches(3.0), Inches(8.3), Inches(3.0),
                summary, font_size=20, color=RGBColor(0xCC, 0xCC, 0xCC))

    add_textbox(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
                "Chris Ruano  |  cruano@g.harvard.edu  |  HEB 91r",
                font_size=18, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
                "Questions?",
                font_size=28, bold=True, color=RED,
                alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "To summarize: SudokuOCR achieves 100% cell-level accuracy on 52 historical census tables. "
        "Over 16,000 cells verified through 26,000 constraint checks, zero errors. "
        "It costs about 8 cents per table -- cheaper and more accurate than manual transcription. "
        "And it's self-verifying: the constraint pass rate tells you quality without any ground truth. "
        "Happy to take questions.")


# ══════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_why_hard(prs)
    slide_04_baseline(prs)
    slide_05_sudoku_insight(prs)
    slide_06_constraints(prs)
    slide_07_architecture(prs)
    slide_08_repair(prs)
    slide_09_results(prs)
    slide_10_cost(prs)
    slide_11_gt_errors(prs)
    slide_12_broader_principle(prs)
    slide_13_future(prs)
    slide_14_thankyou(prs)

    prs.save(str(OUT))
    print(f"Saved {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
